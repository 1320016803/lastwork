"""完整测试：生成 docx/pdf/txt/xlsx/html/ppt 6 种示例文件，全部转换为 md。
运行方式：在项目目录执行 `python _test_full.py`"""
from __future__ import annotations

import json
import sys
from pathlib import Path

APP_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(APP_DIR))

from converters.converter_dispatcher import convert_file, supported_extensions

SAMPLES = APP_DIR / "_sample_inputs"
OUTPUTS = APP_DIR / "_sample_outputs"
SAMPLES.mkdir(exist_ok=True)
OUTPUTS.mkdir(exist_ok=True)


def make_txt() -> Path:
    p = SAMPLES / "示例文本.txt"
    p.write_text("这是一个示例 TXT 文件。\n用于测试中文内容的读取与转换。\n\n第二段落。\n", encoding="utf-8")
    return p


def make_docx() -> Path:
    from docx import Document
    doc = Document()
    doc.add_heading("示例 Word 文档", level=1)
    doc.add_paragraph("这是第一段落，包含普通文本。")
    doc.add_heading("二级标题 · 列表示例", level=2)
    doc.add_paragraph("苹果").paragraph_format.style = doc.styles["Normal"]
    doc.add_paragraph("香蕉")
    table = doc.add_table(rows=2, cols=3)
    for i, v in enumerate(["姓名", "年龄", "城市"]):
        table.rows[0].cells[i].text = v
    for i, v in enumerate(["张三", "25", "北京"]):
        table.rows[1].cells[i].text = v
    p = SAMPLES / "示例文档.docx"
    doc.save(str(p))
    return p


def make_xlsx() -> Path:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "成员表"
    ws.append(["姓名", "年龄", "城市", "分数"])
    ws.append(["张三", 25, "北京", 95.5])
    ws.append(["李四", 30, "上海", 88])
    ws.append(["王五", 27, "广州", 76])
    ws2 = wb.create_sheet("备注")
    ws2.append(["这是第二个工作表的内容"])
    p = SAMPLES / "示例表格.xlsx"
    wb.save(str(p))
    return p


def make_pdf() -> Path:
    import fitz
    p = SAMPLES / "示例PDF.pdf"
    pdf_doc = fitz.open()
    page = pdf_doc.new_page()
    page.insert_text((50, 50), "PDF 示例文档", fontsize=16)
    page.insert_text((50, 90), "这是从 PDF 提取的示例文本。", fontsize=11)
    page.insert_text((50, 120), "用于验证 PyMuPDF 的文本提取能力。", fontsize=11)
    pdf_doc.save(str(p))
    pdf_doc.close()
    return p


def make_html() -> Path:
    p = SAMPLES / "示例网页.html"
    html = """<!DOCTYPE html>
<html>
<head><meta charset="utf-8"><title>示例网页</title></head>
<body>
<h1>这是标题</h1>
<p>这是一段 <strong>加粗</strong> 与 <em>斜体</em> 的文本。</p>
<ul>
  <li>列表项一</li>
  <li>列表项二</li>
  <li>列表项三</li>
</ul>
<table>
  <tr><th>姓名</th><th>年龄</th></tr>
  <tr><td>张三</td><td>25</td></tr>
  <tr><td>李四</td><td>30</td></tr>
</table>
</body>
</html>
"""
    p.write_text(html, encoding="utf-8")
    return p


def make_ppt() -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # 幻灯片 1：标题+内容布局
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "演示文稿示例"
    slide.placeholders[1].text = "副标题：PPT to Markdown"

    # 幻灯片 2：列表布局
    bullet_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_layout)
    slide2.shapes.title.text = "第二页 · 要点"
    body = slide2.placeholders[1].text_frame
    body.text = "要点一"
    body.add_paragraph().text = "要点二"
    body.add_paragraph().text = "要点三"

    # 幻灯片 3：表格
    table_layout = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(table_layout)
    slide3.shapes.title.text = "第三页 · 表格示例"
    rows, cols = 3, 3
    left, top, width, height = Inches(1), Inches(2), Inches(6), Inches(2)
    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for i, v in enumerate(["列 A", "列 B", "列 C"]):
        table.cell(0, i).text = v
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = f"({r},{c})"

    p = SAMPLES / "示例PPT.pptx"
    prs.save(str(p))
    return p


def run() -> int:
    print("支持的扩展名:", supported_extensions())
    print()

    makers = {
        "TXT": make_txt,
        "DOCX": make_docx,
        "XLSX": make_xlsx,
        "PDF": make_pdf,
        "HTML": make_html,
        "PPTX": make_ppt,
    }

    results: list[str] = []
    errors: list[str] = []

    for name, mk in makers.items():
        try:
            src = mk()
        except Exception as exc:
            msg = f"[FAIL] {name} 生成示例文件失败：{exc}"
            print(msg)
            errors.append(msg)
            continue
        try:
            out = convert_file(src, OUTPUTS)
            size = Path(out).stat().st_size
            content = Path(out).read_text(encoding="utf-8").strip()
            lines = content.splitlines()
            preview = "\n".join("   | " + ln for ln in lines[:8])
            print(f"[OK] {name}  {src.name}  →  {Path(out).name}  ({size} 字节)")
            print(preview)
            if len(lines) > 8:
                print(f"    ... 省略 {len(lines) - 8} 行 ...")
            print()
            results.append(name)
        except Exception as exc:
            msg = f"[FAIL] {name}: {exc}"
            print(msg)
            errors.append(msg)
            print()

    print("=" * 50)
    print(f"成功：{len(results)} 项  ({', '.join(results)})")
    print(f"失败：{len(errors)} 项")
    for e in errors:
        print("  -", e)

    # 保存一个总览 json 便于调试
    summary = {
        "supported_extensions": supported_extensions(),
        "success": results,
        "errors": errors,
    }
    (OUTPUTS / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")

    return 0 if not errors else 1


if __name__ == "__main__":
    sys.exit(run())
