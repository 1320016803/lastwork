from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .converter_base import BaseConverter


class PptConverter(BaseConverter):
    SUPPORTED_EXT = (".pptx", ".ppt")

    def _to_markdown(self, src: Path) -> str:
        prs = Presentation(str(src))
        sections: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = f"## 第 {slide_idx} 页\n"
            lines: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE and hasattr(shape, "table"):
                        table_lines = self._render_table(shape.table)
                        if table_lines:
                            lines.append("")
                            lines.extend(table_lines)
                            lines.append("")
                    continue

                tf = shape.text_frame
                for para_idx, paragraph in enumerate(tf.paragraphs):
                    text = self._render_paragraph(paragraph)
                    if not text.strip():
                        lines.append("")
                        continue

                    level = paragraph.level or 0
                    if para_idx == 0 and level == 0 and self._looks_like_title(text):
                        lines.append(f"### {text}")
                    elif level == 0:
                        lines.append(text)
                    else:
                        indent = "  " * (level - 1)
                        lines.append(f"{indent}- {text}")

            page_content = "\n".join(self._collapse_blank_lines(lines)).strip()
            if page_content:
                sections.append(title + "\n" + page_content + "\n")
            else:
                sections.append(title + "\n（空页）\n")

        return "\n".join(sections).strip() + "\n"

    def _render_paragraph(self, paragraph) -> str:
        parts: list[str] = []
        for run in paragraph.runs:
            text = run.text
            if not text:
                continue
            # python-pptx 的粗体/斜体/下划线都放在 run.font 下
            font = run.font
            is_bold = getattr(font, "bold", None) is True
            is_italic = getattr(font, "italic", None) is True
            is_underline = getattr(font, "underline", None) not in (None, False, 0)

            if is_bold and is_italic:
                text = f"***{text}***"
            elif is_bold:
                text = f"**{text}**"
            elif is_italic:
                text = f"*{text}*"
            if is_underline:
                text = f"<u>{text}</u>"
            parts.append(text)
        return "".join(parts).strip()

    def _render_table(self, table) -> list[str]:
        if not table.rows:
            return []

        rendered: list[list[str]] = []
        for row in table.rows:
            rendered.append([self._cell_text(cell) for cell in row.cells])

        n_cols = max(len(row) for row in rendered)
        normalized = [row + [""] * (n_cols - len(row)) for row in rendered]

        header_row = [self._escape(cell) for cell in normalized[0]]
        md_lines = [
            "| " + " | ".join(header_row) + " |",
            "| " + " | ".join(["---"] * n_cols) + " |",
        ]
        for row in normalized[1:]:
            md_lines.append("| " + " | ".join(self._escape(cell) for cell in row) + " |")
        return md_lines

    def _cell_text(self, cell) -> str:
        return " ".join(p.text for p in cell.text_frame.paragraphs).strip()

    def _escape(self, text: str) -> str:
        return text.replace("\n", " ").replace("|", "\\|")

    def _looks_like_title(self, text: str) -> bool:
        stripped = text.strip()
        if not stripped:
            return False
        if len(stripped) > 40:
            return False
        return not stripped.endswith((".", "。", "，", ",", "；", ";", "！", "!", "？", "?"))

    def _collapse_blank_lines(self, lines: list[str]) -> list[str]:
        result: list[str] = []
        blank_count = 0
        for line in lines:
            if line.strip() == "":
                blank_count += 1
                if blank_count <= 2:
                    result.append("")
            else:
                blank_count = 0
                result.append(line)
        return result
